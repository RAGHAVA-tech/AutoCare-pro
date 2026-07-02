from pptx import Presentation

slides = [
    {"title": "Automotive Service Booking — Pitch Deck","bullets": [
        "Smart online booking and shop management for auto service centers",
        "Reduces no-shows, optimizes technician schedules, and increases revenue"
    ]},
    {"title": "Problem","bullets": [
        "Customers face long wait times and manual booking hassles",
        "Shops suffer from no-shows, poor capacity utilization, and administrative overhead"
    ]},
    {"title": "Solution","bullets": [
        "AI-assisted booking flow with confirmations, reminders, and technician matching",
        "Admin dashboard for real-time slots, rescheduling, and reporting"
    ]},
    {"title": "Product Highlights","bullets": [
        "Customer booking, reschedule, cancellation and confirmations",
        "Technician assignment, calendar sync, SMS/email reminders",
        "Payments and invoicing integration, reporting and analytics"
    ]},
    {"title": "Market Opportunity","bullets": [
        "Large, fragmented local auto service market with recurring demand",
        "High TAM for SaaS scheduling + payments in automotive aftercare"
    ]},
    {"title": "Business Model","bullets": [
        "SaaS subscription per location or per-user pricing",
        "Add-ons: payments, SMS credits, premium support, white-label"
    ]},
    {"title": "Go-to-Market","bullets": [
        "Direct sales to regional shop chains and franchise groups",
        "Partnerships with POS, parts suppliers, and local aggregators"
    ]},
    {"title": "Competition","bullets": [
        "Traditional booking tools and generic scheduling apps",
        "Differentiator: AI-driven matching, automotive-specific workflows"
    ]},
    {"title": "Roadmap","bullets": [
        "Q3: Payment & calendar integrations; Q4: multi-location management",
        "Next: Marketplace integrations and advanced analytics"
    ]},
    {"title": "Team","bullets": [
        "Product, automotive ops, and engineering with SaaS experience",
        "Advisors from automotive retail and payments"
    ]},
    {"title": "Financial Snapshot","bullets": [
        "Subscription revenue model; path to profitability with low CAC",
        "3-year projections and break-even by year 2 (example figures)"
    ]},
    {"title": "Ask","bullets": [
        "Seeking X funding to accelerate integrations, sales, and growth",
        "Use of funds: engineering, GTM, and customer success"
    ]}
]

prs = Presentation()
# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Automotive Service Booking"
slide.placeholders[1].text = "AI-powered booking and shop management"

# Content slides
for s in slides:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s['title']
    tf = slide.placeholders[1].text_frame
    bullets = s['bullets']
    if bullets:
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1

output = 'pitch_deck.pptx'
prs.save(output)
print(f"Created {output}")
